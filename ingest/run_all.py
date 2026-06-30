"""
Cider Institute — Knowledge Vault Ingest Pipeline
Converts course materials (PDF, DOCX, PPTX, XLSX, MP4) to markdown files organized by topic.
"""

import os
import sys
from pathlib import Path

# ── Configuration ──
SOURCE_DIR = Path("c:/Users/jfpru/Desktop/proyectos/cider-institute/Jacobo Prudot - Cider Institute")
VAULT_DIR = Path("c:/Users/jfpru/Desktop/proyectos/cider-knowledge/vault")

# Topic mapping: source filename keyword → vault subdirectory
TOPIC_MAP = {
    # Foundation Textbook → fermentation
    "textbook": "fermentation",
    "foundation": "fermentation",

    # Lab Manual → lab-testing
    "lab manual": "lab-testing",
    "laboratory": "lab-testing",

    # Sensory slides → sensory-analysis
    "sensory": "sensory-analysis",

    # Aroma slides → aroma-chemistry
    "aroma": "aroma-chemistry",

    # Equipment + Facility → facility-operations
    "equipment": "facility-operations",
    "facility": "facility-operations",
    "sop": "facility-operations",
    "hedonic": "sensory-analysis",
    "score sheet": "sensory-analysis",

    # Perry videos → perry-production
    "perry": "perry-production",

    # Foundation video → fermentation
    "foundation_downstream": "fermentation",

    # Tech barrels → facility-operations
    "barrels": "facility-operations",
}


def find_topic(filename: str) -> str:
    """Determine which vault topic a source file belongs to."""
    name_lower = filename.lower()
    for keyword, topic in TOPIC_MAP.items():
        if keyword in name_lower:
            return topic
    return "fermentation"  # default


def extract_pdf(filepath: Path) -> list[dict]:
    """Extract text pages from PDF using pypdf."""
    from pypdf import PdfReader
    reader = PdfReader(str(filepath))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"num": i + 1, "text": text.strip()})
    return pages


def extract_docx(filepath: Path) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document
    doc = Document(str(filepath))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            if para.style.name.startswith("Heading"):
                level = para.style.name.split()[-1]
                try:
                    level = int(level)
                    prefix = "#" * min(level, 4)
                except ValueError:
                    prefix = "##"
                parts.append(f"\n{prefix} {para.text.strip()}\n")
            else:
                parts.append(para.text.strip())

    # Extract tables
    for ti, table in enumerate(doc.tables):
        parts.append(f"\n### Table {ti + 1}\n")
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            header = rows[0]
            sep = "|" + "|".join([" --- " for _ in table.columns]) + "|"
            parts.append(header)
            parts.append(sep)
            parts.extend(rows[1:])

    return "\n\n".join(parts)


def extract_pptx(filepath: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation
    prs = Presentation(str(filepath))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            parts.append(f"\n### Slide {i + 1}\n")
            parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def extract_xlsx(filepath: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook
    wb = load_workbook(str(filepath), data_only=True)
    parts = []
    for sname in wb.sheetnames:
        ws = wb[sname]
        parts.append(f"\n## Sheet: {sname}\n")
        rows_list = []
        for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, 200), values_only=True):
            vals = [str(v) if v is not None else "" for v in row]
            if any(v for v in vals):
                rows_list.append("| " + " | ".join(vals[:15]) + " |")  # limit columns
        if rows_list:
            parts.extend(rows_list)
    return "\n".join(parts)


def transcribe_video(filepath: Path, model_size: str = "tiny") -> str:
    """Transcribe video using OpenAI Whisper.

    Requires: pip install openai-whisper
    Whisper reads video files directly — no ffmpeg extraction needed.

    Model sizes: tiny (fastest, ~1GB RAM), base (~1.5GB), small (~2.5GB), medium (~5GB)
    Speed on ARM CPU (tiny): ~10x realtime for English
    """
    import time
    import whisper

    file_size_mb = filepath.stat().st_size / (1024 * 1024)
    print(f"   Loading whisper model '{model_size}'...")
    t0 = time.time()
    model = whisper.load_model(model_size)
    print(f"   Model loaded in {time.time() - t0:.0f}s")

    print(f"   Transcribing ({file_size_mb:.0f} MB video)...")
    t0 = time.time()
    result = model.transcribe(str(filepath), language="en")
    elapsed = time.time() - t0

    words = len(result["text"].split())
    print(f"   Done in {elapsed:.0f}s — {words} words ({words / elapsed:.0f} words/sec)")

    return result["text"]


def save_to_vault(topic: str, title: str, content: str):
    """Save content to the appropriate vault topic directory."""
    topic_dir = VAULT_DIR / topic
    topic_dir.mkdir(parents=True, exist_ok=True)

    filename = title.lower().replace(" ", "-").replace("/", "-")[:80] + ".md"
    filepath = topic_dir / filename

    header = f"# {title}\n\n"
    header += f"> Source: Cider Institute of North America course materials\n\n"

    filepath.write_text(header + content, encoding="utf-8")
    print(f"  ✓ Saved: {topic}/{filename}")
    return filepath


def process_file(filepath: Path) -> list:
    """Process a single source file and return list of created vault files."""
    ext = filepath.suffix.lower()
    name = filepath.stem
    topic = find_topic(name)
    created = []

    print(f"\n📄 Processing: {filepath.name} → {topic}/")

    try:
        if ext == ".pdf":
            pages = extract_pdf(filepath)
            print(f"   {len(pages)} pages extracted")
            # For large PDFs, group pages into logical chunks
            if len(pages) <= 5:
                # Save each page or small groups
                content = "\n\n".join(p["text"] for p in pages)
                fp = save_to_vault(topic, name, content)
                created.append(fp)
            else:
                # Split into sections (every ~10 pages or by headings)
                chunk = []
                chunk_num = 1
                for p in pages:
                    chunk.append(p["text"])
                    if len(chunk) >= 10:
                        fp = save_to_vault(topic, f"{name} part {chunk_num}", "\n\n".join(chunk))
                        created.append(fp)
                        chunk = []
                        chunk_num += 1
                if chunk:
                    fp = save_to_vault(topic, f"{name} part {chunk_num}", "\n\n".join(chunk))
                    created.append(fp)

        elif ext == ".docx":
            content = extract_docx(filepath)
            # Split large docs by H2 headings
            sections = content.split("\n## ")
            if len(sections) > 1:
                # First section is intro
                if sections[0].strip():
                    fp = save_to_vault(topic, name, sections[0].strip())
                    created.append(fp)
                for sec in sections[1:]:
                    lines = sec.strip().split("\n", 1)
                    sec_title = lines[0].strip()
                    sec_body = lines[1] if len(lines) > 1 else ""
                    fp = save_to_vault(topic, f"{name} — {sec_title}", f"## {sec_title}\n{sec_body}")
                    created.append(fp)
            else:
                fp = save_to_vault(topic, name, content)
                created.append(fp)

        elif ext == ".pptx":
            content = extract_pptx(filepath)
            fp = save_to_vault(topic, name, content)
            created.append(fp)

        elif ext == ".xlsx":
            content = extract_xlsx(filepath)
            fp = save_to_vault(topic, name, content)
            created.append(fp)

        elif ext == ".mp4":
            print("   🎬 Transcribing video (this may take several minutes)...")
            transcript = transcribe_video(filepath)
            fp = save_to_vault(topic, f"{name} transcript", transcript)
            created.append(fp)

        else:
            print(f"   ⚠ Unsupported format: {ext}")

    except Exception as e:
        print(f"   ✗ Error: {e}")

    return created


def main():
    """Run the full ingest pipeline."""
    print("🍎 Cider Institute — Knowledge Vault Ingest")
    print("=" * 50)

    # Find all source files
    source_files = sorted(
        [f for f in SOURCE_DIR.iterdir() if f.suffix.lower() in {".pdf", ".docx", ".pptx", ".xlsx", ".mp4"}],
        key=lambda f: f.name
    )

    print(f"\nFound {len(source_files)} source files in:\n  {SOURCE_DIR}")
    print(f"Output:\n  {VAULT_DIR}\n")

    total_created = 0
    for filepath in source_files:
        created = process_file(filepath)
        total_created += len(created)

    print(f"\n{'=' * 50}")
    print(f"✅ Ingest complete. {total_created} markdown files created.")
    print(f"   Vault: {VAULT_DIR}")


if __name__ == "__main__":
    main()
